from helper import print_formatted, print_delta, fix_name, print_comment
from helper import REVERSE_METRICS, filter_data_by_kpi, get_metric_type
from dates import yesterday, sdlw, get_previous_week_start_date_end_date, timedelta
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement


def add_kpi_heading(slide, asset):
    main_heading_text_left = Inches(0.15)
    main_heading_text_top = Inches(0.15)
    main_heading_text_width = Inches(10)
    main_heading_text_height = Inches(1)
    main_heading_textbox = slide.shapes.add_textbox(
        main_heading_text_left, main_heading_text_top, main_heading_text_width, main_heading_text_height
    )

    main_heading_text_frame = main_heading_textbox.text_frame
    main_heading_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = main_heading_text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f'{fix_name(asset)} Key Metrics Performance'
    font = run.font
    font.name = 'Poppins'
    font.size = Pt(18)
    font.bold = True
    font.color.rgb = RGBColor(70, 177, 255)


def add_date(slide, period, weekday):
    date_shape_left = Inches(10.5)
    date_shape_top = Inches(0.35)
    date_shape_width = Cm(5.89)
    date_shape_height = Cm(1.66)
    date_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, date_shape_left, date_shape_top, date_shape_width, date_shape_height
    )
    fill = date_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(176, 231, 198)
    shadow = date_shape.shadow
    shadow.inherit = False
    line = date_shape.line
    line.fill.background()

    date_text_left = Inches(date_shape_left.inches - 0.05)
    date_text_top = date_shape_top
    date_text_width = date_shape_width
    date_text_height = date_shape_height
    date_textbox = slide.shapes.add_textbox(
        date_text_left, date_text_top, date_text_width, date_text_height
    )

    date_text_frame = date_textbox.text_frame
    date_text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = date_text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()

    if period == 'daily':
        date_text = f'Date - {yesterday.strftime("%d/%b/%y")}'
    elif period == 'weekly':
        week_start, week_end = get_previous_week_start_date_end_date(weekday=weekday)
        date_text = f'Week {week_start.strftime("%U")}'
    else:
        date_text = f'Unknown period - {period} '

    run.text = date_text
    font = run.font
    font.name = 'Poppins'
    font.size = Pt(18)
    font.bold = True

    date_subtext_left = date_text_left
    date_subtext_top = date_text_top
    date_subtext_width = date_text_width
    date_subtext_height = date_text_height
    date_subtextbox = slide.shapes.add_textbox(
        date_subtext_left, date_subtext_top, date_subtext_width, date_subtext_height
    )

    date_subtext_frame = date_subtextbox.text_frame
    date_subtext_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = date_subtext_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()

    if period == 'daily':
        date_subtext = ''
    elif period == 'weekly':
        week_start, week_end = get_previous_week_start_date_end_date(weekday=weekday)
        date_subtext = f'{week_start.strftime("%d-%b-%y")} – {week_end.strftime("%d-%b-%y")}'
    else:
        date_subtext = f'Unknown period - {period} '

    run.text = date_subtext
    font = run.font
    font.name = 'Poppins'
    font.size = Pt(11)
    font.bold = True


def add_card_shape(kpi_slide, left, top, width, height):
    card_shape = kpi_slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )

    card_shape.adjustments[0] = 0.05

    fill = card_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    card_shape.shadow.inherit = True

    card_shape.line.fill.background()


def add_metric_name(kpi_slide, left, top, width, height, metric):
    metric_name_textbox = kpi_slide.shapes.add_textbox(
        left, top, width, height
    )

    metric_name_text_frame = metric_name_textbox.text_frame
    metric_name_text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = metric_name_text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()

    metric_name_text = fix_name(metric)

    run.text = metric_name_text
    font = run.font
    font.name = 'Poppins'
    font.size = Pt(10)
    font.bold = True
    font.color.rgb = RGBColor(166, 166, 166)


def add_metric_value(kpi_slide, left, top, width, height, y, metric):
    metric_value_textbox = kpi_slide.shapes.add_textbox(
        left, top, width, height
    )

    metric_value_text_frame = metric_value_textbox.text_frame
    metric_value_text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = metric_value_text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()

    metric_value_text = print_formatted(y, metric)

    run.text = metric_value_text
    font = run.font
    font.name = 'Poppins'
    font.size = Pt(18)
    font.bold = True


def add_delta(kpi_slide, left, top, width, height, period, metric, y, y_prev):
    delta_textbox = kpi_slide.shapes.add_textbox(
        left, top, width, height
    )

    delta_text_frame = delta_textbox.text_frame
    delta_text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = delta_text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()

    if period == 'daily':
        delta_text = f"{print_delta(now=y, prev=y_prev)} vs SDLW"
    elif period == 'weekly':
        delta_text = f"WoW {print_delta(now=y, prev=y_prev)}"
    else:
        delta_text = f"Invalid period {print_delta(now=y, prev=y_prev)}"

    run.text = delta_text
    font = run.font
    font.name = 'Arial'
    font.size = Pt(11)
    font.bold = True
    if y > y_prev:
        if metric.lower() in REVERSE_METRICS:
            font.color.rgb = RGBColor(229, 88, 1)
        else:
            font.color.rgb = RGBColor(0, 176, 80)
    else:
        if metric.lower() in REVERSE_METRICS:
            font.color.rgb = RGBColor(0, 176, 80)
        else:
            font.color.rgb = RGBColor(229, 88, 1)


def add_divider(kpi_slide, left, top, width, height, metric, y, y_prev):
    divider_shape = kpi_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )

    fill = divider_shape.fill
    fill.solid()
    if y > y_prev:
        if metric.lower() in REVERSE_METRICS:
            fill.fore_color.rgb = RGBColor(229, 88, 1)
        else:
            fill.fore_color.rgb = RGBColor(0, 176, 80)
    else:
        if metric.lower() in REVERSE_METRICS:
            fill.fore_color.rgb = RGBColor(0, 176, 80)
        else:
            fill.fore_color.rgb = RGBColor(229, 88, 1)

    shadow = divider_shape.shadow
    shadow.inherit = False
    line = divider_shape.line
    line.fill.background()


def add_kpi_chart(kpi_slide, left, top, width, height, kpi, asset_df):
    filtered_df = filter_data_by_kpi(asset_df, kpi)

    y = filtered_df['y'].iloc[0]
    y_prev = filtered_df['y_prev'].iloc[0]

    chart_data = ChartData()

    chart_data.categories = filtered_df['xaxis_data'].iloc[0]
    chart_data.add_series('Series 1', filtered_df['yaxis_data'].iloc[0])

    chart = kpi_slide.shapes.add_chart(
        XL_CHART_TYPE.AREA, left, top, width, height, chart_data
    ).chart
    chart.has_title = False
    chart.has_legend = False

    xaxis = chart.category_axis
    xaxis.has_major_gridlines = False
    xaxis.tick_labels.font.size = Pt(8)
    # xaxis.visible = False

    yaxis = chart.value_axis
    yaxis.minimum_scale = 0
    yaxis.has_major_gridlines = False
    yaxis.tick_labels.font.bold = True
    yaxis.tick_labels.font.size = Pt(8)
    yaxis.tick_labels.font.color.rgb = RGBColor(82, 99, 142)

    plot = chart.plots[0]
    series = plot.series[0]
    fill = series.format.fill
    fill.solid()
    if y > y_prev:
        if kpi.metric.lower() in REVERSE_METRICS:
            fill.fore_color.rgb = RGBColor(229, 88, 1)
        else:
            fill.fore_color.rgb = RGBColor(0, 176, 80)
    else:
        if kpi.metric.lower() in REVERSE_METRICS:
            fill.fore_color.rgb = RGBColor(0, 176, 80)
        else:
            fill.fore_color.rgb = RGBColor(229, 88, 1)
    solidFill = fill.fore_color._xFill
    alpha = OxmlElement('a:alpha')
    alpha.set('val', '50196')
    solidFill.srgbClr.append(alpha)

    line = series.format.line
    if y > y_prev:
        if kpi.metric.lower() in REVERSE_METRICS:
            line.color.rgb = RGBColor(229, 88, 1)
        else:
            line.color.rgb = RGBColor(0, 72, 33)
    else:
        if kpi.metric.lower() in REVERSE_METRICS:
            line.color.rgb = RGBColor(0, 72, 33)
        else:
            line.color.rgb = RGBColor(229, 88, 1)

    metric_type = get_metric_type(kpi.metric)

    if metric_type == 'revenue':
        yaxis.tick_labels.number_format = '[<1000]"$"0.0;[<999950]"$"0.0,"K";[<999950000]"$"0.0,,"M";"$"0.0,,,"B"'
    elif metric_type == 'aov':
        yaxis.tick_labels.number_format = '"$"0'
    elif metric_type == 'rate':
        yaxis.tick_labels.number_format = '0%'
    else:
        yaxis.tick_labels.number_format = '[<1000]0.0;[<999950]0.0,"K";[<999950000]0.0,,"M";0.0,,,"B"'


def add_kpi_card(kpi_slide, kpi, period, asset_df, left, top):
    print(f"Adding kpi card for {kpi.metric}")
    filtered_df = filter_data_by_kpi(asset_df, kpi)

    if filtered_df.empty:
        print(f"KPI - {kpi.metric} not available")
        return

    y = filtered_df['y'].iloc[0]
    y_prev = filtered_df['y_prev'].iloc[0]

    width = Cm(4.64)
    height = Cm(5.71)

    add_card_shape(kpi_slide, left, top, width, height)

    add_metric_name(
        kpi_slide,
        left=Inches(left.inches + 0.05),
        top=Inches(top.inches + 0.15),
        width=width,
        height=height,
        metric=kpi.metric
    )

    add_metric_value(
        kpi_slide,
        left=Inches(left.inches + 0.05),
        top=Inches(top.inches + 0.35),
        width=width, height=height,
        y=filtered_df['y'].iloc[0],
        metric=kpi.metric
    )

    add_delta(
        kpi_slide,
        left,
        top=Inches(top.inches + 0.7),
        width=width,
        height=height,
        period=period,
        metric=kpi.metric,
        y=y,
        y_prev=y_prev
    )

    add_divider(
        kpi_slide,
        left=Inches(left.inches + width.inches / 2 - Cm(3.24).inches / 2),
        top=Inches(top.inches + 1.05),
        width=Cm(3.24),
        height=Cm(0.12),
        metric=kpi.metric,
        y=y,
        y_prev=y_prev
    )

    add_kpi_chart(
        kpi_slide,
        left,
        top=Inches(top.inches + 1.2), width=Cm(4.8),
        height=Cm(2.5),
        kpi=kpi,
        asset_df=asset_df
    )


def add_comments(kpi_slide, kpi_list, asset_df):
    comments_text_left = Inches(7.5)
    comments_text_top = Inches(1.25)
    comments_text_width = Cm(13.81)
    comments_text_height = Cm(12.87)
    comments_textbox = kpi_slide.shapes.add_textbox(
        comments_text_left, comments_text_top, comments_text_width, comments_text_height
    )

    comments_text_frame = comments_textbox.text_frame
    comments_text_frame.vertical_anchor = MSO_ANCHOR.TOP

    kpi_slide.shapes.add_picture(
        'shadow.png',
        left=Inches(comments_text_left.inches - 0.5),
        top=comments_text_top,
        width=Cm(0.94),
        height=Cm(15.14)
    )

    for kpi in kpi_list:
        filtered_df = filter_data_by_kpi(asset_df, kpi)
        if filtered_df.empty:
            print(f"KPI - {kpi.metric} not available")
            return

        p = comments_text_frame.add_paragraph()
        run = p.add_run()

        comments_text = print_comment(filtered_df.iloc[0, :])

        run.text = comments_text
        font = run.font
        font.name = 'Poppins'
        font.size = Pt(12)
        font.bold = True
        font.color.rgb = RGBColor(0, 32, 96)

        p = comments_text_frame.add_paragraph()
        p.level = 1
        run = p.add_run()

        comments_text = 'Point 1'

        run.text = comments_text
        font = run.font
        font.name = 'Poppins'
        font.size = Pt(12)
        font.color.rgb = RGBColor(0, 32, 96)

        comments_text_frame.add_paragraph()


def add_kpi_cards(kpi_slide, period, asset_df, kpi_list):
    card_locations = [
        (Inches(0.5), Inches(1.5)),
        (Inches(2.5), Inches(1.5)),
        (Inches(4.5), Inches(1.5)),
        (Inches(0.5), Inches(4.0)),
        (Inches(2.5), Inches(4.0)),
        (Inches(4.5), Inches(4.0)),
    ]

    for kpi, (left, top) in zip(kpi_list, card_locations):
        add_kpi_card(kpi_slide, kpi, period, asset_df, left, top)

    add_comments(kpi_slide, kpi_list, asset_df)


def add_kpi_slide(ppt, period, asset, asset_df, kpi_list):
    blank_slide_layout = ppt.slide_layouts[6]
    kpi_slide = ppt.slides.add_slide(blank_slide_layout)
    add_kpi_heading(kpi_slide, asset)

    if period == 'weekly':
        weekday = asset_df['weekday'].mode().iloc[0]
        if not asset_df[asset_df['weekday'] != weekday].empty:
            raise ValueError(f"Weekdays are different for {asset_df[asset_df['weekday'] != weekday]}")
    else:
        weekday = None

    add_date(kpi_slide, period, weekday)
    add_kpi_cards(kpi_slide, period, asset_df, kpi_list)
