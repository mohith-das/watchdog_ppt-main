from rca import print_revenue_rca
from helper import print_formatted, print_delta, fix_name, always_include_data_source
from dates import yesterday, sdlw, get_previous_week_start_date_end_date
from graph import get_graph
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
import plotly.io as pio


def add_rca_slide(ppt, period, asset_df):
    blank_slide_layout = ppt.slide_layouts[6]
    rca_slide = ppt.slides.add_slide(blank_slide_layout)

    left = Cm(0.5)
    top = Cm(0.5)
    width = Cm(6.81)
    height = Cm(7.98)

    dimension_textbox = rca_slide.shapes.add_textbox(
        left, top, width, height
    )

    dimension_text_frame = dimension_textbox.text_frame
    dimension_text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = dimension_text_frame.paragraphs[0]
    print_revenue_rca(asset_df, p)
